"""
filename: report.py
purpose: Build a PowerPoint report from the analysis CSV and charts produced
         by the pipeline. Slides include title, a summary table, the three
         charts (bar, deviation, pie), and a conclusion slide highlighting
         key findings.
author: Tuna Safak
date: 2026-06-14
decision_rationale: A PowerPoint slide deck provides a familiar format for
  presenting results to stakeholders. We create a concise report including
  the summary table and visual charts so the lab results are both human- and
  presentation-ready.

This script uses `pandas` to read `output/analysis_results.csv`, `python-pptx`
to construct slides, and saves the presentation to `output/blood_test_report.pptx`.
Detailed status prints and try/except blocks help track progress and errors.
"""

import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def ensure_output_dir(output_dir: str):
    """Create the output directory if it does not exist."""
    if not os.path.exists(output_dir):
        print(f"Creating output directory: {output_dir}")
        try:
            os.makedirs(output_dir, exist_ok=True)
        except Exception as e:
            print(f"ERROR: Could not create output directory: {e}")
            raise


def add_title_slide(prs, title_text: str, subtitle_text: str = ''):
    """Add a title slide (layout 0) to the presentation."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text


def add_table_slide(prs, df: pd.DataFrame, title: str):
    """Add a slide with a summary table showing marker, value, reference, risk."""
    slide_layout = prs.slide_layouts[5]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    rows = df.shape[0] + 1  # header + rows
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)  # marker
    table.columns[1].width = Inches(1.5)  # value
    table.columns[2].width = Inches(3.0)  # reference
    table.columns[3].width = Inches(1.5)  # risk_level

    # Header
    headers = ['Marker', 'Value', 'Reference', 'Risk Level']
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True

    # Rows
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        table.cell(i, 0).text = str(row.get('marker', ''))
        table.cell(i, 1).text = str(row.get('value', ''))
        table.cell(i, 2).text = str(row.get('reference', ''))
        table.cell(i, 3).text = str(row.get('risk_level', ''))


def add_image_slide(prs, image_path: str, title: str):
    """Add a slide with a single image centered in the content area."""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    try:
        slide.shapes.add_picture(image_path, left, top, width=width)
    except Exception as e:
        print(f"Warning: could not insert image {image_path}: {e}")


def add_conclusion_slide(prs, conclusions: list):
    """Add a conclusion slide with bullet points listing key findings."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Conclusion'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(conclusions):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    analysis_csv = os.path.join(script_dir, 'output', 'analysis_results.csv')
    output_dir = os.path.join(script_dir, 'output')
    pptx_path = os.path.join(output_dir, 'blood_test_report.pptx')

    print('Report generator started.')

    # Ensure output directory exists
    try:
        ensure_output_dir(output_dir)
    except Exception:
        sys.exit(1)

    # Read analysis CSV
    try:
        print(f"Reading analysis CSV: {analysis_csv}")
        df = pd.read_csv(analysis_csv)
    except Exception as e:
        print(f"ERROR: Could not read analysis CSV: {e}")
        sys.exit(1)

    # Start presentation
    try:
        prs = Presentation()
    except Exception as e:
        print(f"ERROR: Could not create Presentation object: {e}")
        sys.exit(1)

    # Slide 1: Title
    add_title_slide(prs, 'Blood Test Analysis Report — April 2021', '')

    # Slide 2: Summary table
    try:
        add_table_slide(prs, df[['marker', 'value', 'reference', 'risk_level']], 'Summary of Results')
    except Exception as e:
        print(f"Warning: could not add summary table: {e}")

    # Slide 3..5: Charts
    charts = [
        ('bar_chart.png', 'Bar Chart — Values by Marker'),
        ('deviation_chart.png', 'Deviation Chart — % Deviation from Center'),
        ('status_pie.png', 'Status Distribution')
    ]
    for filename, title in charts:
        img_path = os.path.join(output_dir, filename)
        add_image_slide(prs, img_path, title)

    # Slide 6: Conclusion
    # Prepare key findings: counts of non-normal risk levels and top deviations
    try:
        counts = df['risk_level'].value_counts().to_dict()
        non_normal = {k: v for k, v in counts.items() if k != 'normal'}
        max_dev_row = df.loc[df['pct_deviation_from_center'].abs().idxmax()]
        conclusions = []
        conclusions.append(f"Total markers evaluated: {len(df)}")
        if non_normal:
            for k, v in non_normal.items():
                conclusions.append(f"{v} marker(s) with risk level: {k}")
        else:
            conclusions.append('All markers are within normal limits based on provided ranges.')

        conclusions.append(f"Largest deviation: {max_dev_row['marker']} ({max_dev_row['pct_deviation_from_center']}% from center)")
        conclusions.append('Recommendations: review borderline markers and consider clinical correlation.')
    except Exception as e:
        conclusions = ['Could not compute key findings due to error: ' + str(e)]

    add_conclusion_slide(prs, conclusions)

    # Save presentation
    try:
        prs.save(pptx_path)
        print(f"Saved PowerPoint: {pptx_path}")
    except Exception as e:
        print(f"ERROR: Could not save PowerPoint: {e}")
        sys.exit(1)

    print('Report generation complete.')


if __name__ == '__main__':
    main()
