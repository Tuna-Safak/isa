"""
Generate a PowerPoint report from structured blood test data.

Input:
    output/bloodtest_results.csv
    optionally output/bloodtest_results.json

Output:
    output/bloodtest_report.pptx
"""

from __future__ import annotations

import os
import json
import re
import sys
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

import pandas as pd
os.environ.setdefault("MPLCONFIGDIR", "/private/tmp/codex-mplconfig")
os.environ.setdefault("XDG_CACHE_HOME", "/private/tmp/codex-cache")
Path(os.environ["MPLCONFIGDIR"]).mkdir(parents=True, exist_ok=True)
Path(os.environ["XDG_CACHE_HOME"]).mkdir(parents=True, exist_ok=True)

try:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except Exception:  # pragma: no cover - optional dependency fallback
    plt = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "output"
CSV_PATH = OUTPUT_DIR / "bloodtest_results.csv"
JSON_PATH = OUTPUT_DIR / "bloodtest_results.json"
METADATA_PATH = OUTPUT_DIR / "bloodtest_metadata.json"
PPTX_PATH = OUTPUT_DIR / "bloodtest_report.pptx"
DOCX_PATH = OUTPUT_DIR / "bloodtest_report.docx"

NAVY = RGBColor(11, 61, 89)
TEAL = RGBColor(21, 140, 146)
MINT = RGBColor(221, 245, 243)
SLATE = RGBColor(73, 84, 100)
SOFT_BG = RGBColor(247, 249, 251)
LINE = RGBColor(216, 224, 232)


@dataclass
class ParsedReference:
    low: Optional[float]
    high: Optional[float]
    mode: str  # "range", "upper", "lower", "unknown"
    raw: str


def ensure_output_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def load_records(csv_path: Path, json_path: Path) -> pd.DataFrame:
    if csv_path.exists() and csv_path.stat().st_size > 0:
        df = pd.read_csv(csv_path, dtype=str)
    elif json_path.exists() and json_path.stat().st_size > 0:
        with json_path.open("r", encoding="utf-8") as handle:
            data = json.load(handle)
        df = pd.DataFrame(data)
    else:
        raise FileNotFoundError("Neither bloodtest_results.csv nor bloodtest_results.json was found.")

    for column in ["marker", "value", "unit", "reference"]:
        if column not in df.columns:
            df[column] = ""

    df["marker"] = df["marker"].fillna("").astype(str)
    df["value"] = df["value"].fillna("").astype(str)
    df["unit"] = df["unit"].fillna("").astype(str)
    df["reference"] = df["reference"].fillna("").astype(str)
    return df


def load_metadata(metadata_path: Path) -> dict:
    if not metadata_path.exists() or metadata_path.stat().st_size == 0:
        return {}
    with metadata_path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    return data if isinstance(data, dict) else {}


def parse_float(text: str) -> Optional[float]:
    if text is None:
        return None
    match = re.search(r"-?\d+(?:[.,]\d+)?", str(text))
    if not match:
        return None
    try:
        return float(match.group(0).replace(",", "."))
    except ValueError:
        return None


def parse_reference(reference: str) -> ParsedReference:
    ref = (reference or "").strip()
    if not ref:
        return ParsedReference(None, None, "unknown", ref)

    ref = re.sub(r"\.(?=\s+\d)", "-", ref)
    ref = re.sub(r"\s+[.\-–]\s+", "-", ref)
    ref_match = re.match(
        r"^\s*(?P<ref>(?:[<>]?\s*\d+[.,]?\d*(?:\s*[-–]\s*\d+[.,]?\d*)?)).*$",
        ref,
    )
    if ref_match:
        ref = ref_match.group("ref")

    range_match = re.search(r"(\d+(?:[.,]\d+)?)\s*[-–]\s*(\d+(?:[.,]\d+)?)", ref)
    if range_match:
        low = parse_float(range_match.group(1))
        high = parse_float(range_match.group(2))
        return ParsedReference(low, high, "range", ref)

    loose_range_match = re.match(r"^\s*(\d+(?:[.,]?\d*)?)\s+(\d+(?:[.,]?\d*)?)\s*$", ref)
    if loose_range_match:
        low = parse_float(loose_range_match.group(1))
        high = parse_float(loose_range_match.group(2))
        if low is not None and high is not None and 0 < low < high and high / max(low, 1e-9) < 20:
            return ParsedReference(low, high, "range", f"{low}-{high}")

    upper_match = re.search(r"^<?\s*(\d+(?:[.,]\d+)?)", ref)
    if ref.startswith("<") and upper_match:
        high = parse_float(upper_match.group(1))
        return ParsedReference(None, high, "upper", ref)

    lower_match = re.search(r"^>?\s*(\d+(?:[.,]\d+)?)", ref)
    if ref.startswith(">") and lower_match:
        low = parse_float(lower_match.group(1))
        return ParsedReference(low, None, "lower", ref)

    return ParsedReference(None, None, "unknown", ref)


def infer_default_reference(marker: str) -> ParsedReference:
    marker_key = (marker or "").strip().lower()

    if "hdl" in marker_key:
        return ParsedReference(40.0, None, "lower", "> 40")
    if "ldl" in marker_key:
        return ParsedReference(None, 130.0, "upper", "< 130")
    if "cholesterin" in marker_key:
        return ParsedReference(None, 200.0, "upper", "< 200")
    if "triglycerid" in marker_key:
        return ParsedReference(None, 150.0, "upper", "< 150")
    if "kreatinin" in marker_key:
        return ParsedReference(None, 1.2, "upper", "< 1.2")
    if "calcium" in marker_key:
        return ParsedReference(2.1, 2.55, "range", "2.1-2.55")
    if "glukose" in marker_key or "glucose" in marker_key:
        return ParsedReference(70.0, 100.0, "range", "70-100")
    if "hba1c" in marker_key:
        return ParsedReference(None, 5.7, "upper", "< 5.7")
    if "gpt" in marker_key or "alat" in marker_key or "alt" in marker_key:
        return ParsedReference(None, 50.0, "upper", "< 50")
    if "gamma" in marker_key or "gt" in marker_key:
        return ParsedReference(None, 50.0, "upper", "< 50")
    if "kalium" in marker_key:
        return ParsedReference(3.5, 5.1, "range", "3.5-5.1")
    if "natrium" in marker_key:
        return ParsedReference(136.0, 145.0, "range", "136-145")
    if "harn" in marker_key and "säure" in marker_key:
        return ParsedReference(None, 7.0, "upper", "< 7.0")

    return ParsedReference(None, None, "unknown", "")


def classify_value(value: float, parsed: ParsedReference) -> str:
    if parsed.mode == "range" and parsed.low is not None and parsed.high is not None:
        if value < parsed.low:
            return "low"
        if value > parsed.high:
            return "high"
        if abs(value - parsed.low) / max(parsed.low, 1e-9) <= 0.1 or abs(parsed.high - value) / max(parsed.high, 1e-9) <= 0.1:
            return "borderline"
        return "normal"

    if parsed.mode == "upper" and parsed.high is not None:
        if value > parsed.high:
            return "high"
        if abs(parsed.high - value) / max(parsed.high, 1e-9) <= 0.1:
            return "borderline"
        return "normal"

    if parsed.mode == "lower" and parsed.low is not None:
        if value < parsed.low:
            return "low"
        if abs(value - parsed.low) / max(parsed.low, 1e-9) <= 0.1:
            return "borderline"
        return "normal"

    return "unknown"


def percent_deviation(value: float, parsed: ParsedReference) -> Optional[float]:
    if parsed.mode == "range" and parsed.low is not None and parsed.high is not None:
        center = (parsed.low + parsed.high) / 2.0
        if center:
            return round(((value - center) / center) * 100.0, 2)
    return None


def enrich_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    records = []
    for _, row in df.iterrows():
        value = parse_float(row["value"])
        parsed = parse_reference(row["reference"])
        if value is not None and parsed.mode == "unknown":
            parsed = infer_default_reference(str(row["marker"]).strip())
        status = "unknown"
        deviation = None

        if value is not None:
            status = classify_value(value, parsed)
            deviation = percent_deviation(value, parsed)

        records.append(
            {
                "marker": str(row["marker"]).strip(),
                "value": value,
                "unit": str(row["unit"]).strip(),
                "reference": parsed.raw,
                "status": status,
                "pct_deviation_from_center": deviation,
                "ref_min": parsed.low,
                "ref_max": parsed.high,
            }
        )

    out = pd.DataFrame(records)
    out = out.drop_duplicates(subset=["marker", "value", "unit", "reference"], keep="first").reset_index(drop=True)
    return out


def color_for_status(status: str) -> str:
    return {
        "normal": "#2e7d32",
        "borderline": "#f9a825",
        "low": "#ef6c00",
        "high": "#c62828",
        "unknown": "#546e7a",
    }.get(status, "#546e7a")


def add_modern_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SOFT_BG
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.38))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.38), Inches(10), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def add_slide_title(slide, title: str, subtitle: Optional[str] = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(8.8), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.56), Inches(1.05), Inches(9.0), Inches(0.35))
        stf = sub_box.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(11)
        stf.paragraphs[0].font.color.rgb = SLATE


def build_bar_chart(df: pd.DataFrame, out_path: Path) -> Optional[Path]:
    if plt is None:
        return None
    usable = df[df["value"].notna()].copy()
    if usable.empty:
        return None

    usable = usable.sort_values(by="value", ascending=False)
    plt.figure(figsize=(12, 7), facecolor="white")
    colors = [color_for_status(str(status)) for status in usable["status"]]
    plt.bar(usable["marker"], usable["value"], color=colors, edgecolor="#2c3e50", linewidth=0.7)
    plt.xticks(rotation=45, ha="right")
    plt.ylabel("Measured value")
    plt.title("Blood Test Values", fontsize=16, fontweight="bold")
    plt.grid(axis="y", alpha=0.18)
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()
    return out_path


def build_deviation_chart(df: pd.DataFrame, out_path: Path) -> Optional[Path]:
    if plt is None:
        return None
    usable = df[df["pct_deviation_from_center"].notna()].copy()
    if usable.empty:
        return None

    usable = usable.sort_values(by="pct_deviation_from_center")
    plt.figure(figsize=(12, 7), facecolor="white")
    colors = [color_for_status(str(status)) for status in usable["status"]]
    plt.barh(usable["marker"], usable["pct_deviation_from_center"], color=colors, edgecolor="#2c3e50", linewidth=0.7)
    plt.axvline(0, color="black", linewidth=1)
    plt.xlabel("% deviation from reference center")
    plt.title("Deviation from Reference Center", fontsize=16, fontweight="bold")
    plt.grid(axis="x", alpha=0.18)
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()
    return out_path


def build_status_pie(df: pd.DataFrame, out_path: Path) -> Optional[Path]:
    if plt is None:
        return None
    counts = Counter(df["status"].fillna("unknown"))
    if not counts:
        return None

    labels = list(counts.keys())
    sizes = [counts[label] for label in labels]
    colors = [color_for_status(label) for label in labels]

    plt.figure(figsize=(8, 6), facecolor="white")
    plt.pie(sizes, labels=labels, colors=colors, autopct="%1.0f%%", startangle=90)
    plt.title("Status Distribution", fontsize=16, fontweight="bold")
    plt.axis("equal")
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()
    return out_path


def add_title_slide(prs: Presentation, df: pd.DataFrame, test_date: Optional[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    teal_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.45), Inches(4.2), Inches(4.9))
    teal_block.fill.solid()
    teal_block.fill.fore_color.rgb = NAVY
    teal_block.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(3.5), Inches(4.7), Inches(1.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MINT
    accent.line.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(6.8), Inches(1.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Doctor's Report"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    sub = tf.add_paragraph()
    sub.text = "Structured OCR results from a German laboratory report"
    sub.font.size = Pt(14)
    sub.font.color.rgb = RGBColor(225, 238, 242)

    stats = slide.shapes.add_textbox(Inches(0.78), Inches(2.75), Inches(3.0), Inches(1.0))
    stf = stats.text_frame
    stf.clear()
    stat_lines = [
        f"{len(df)} markers extracted",
        f"{int((df['status'].isin(['low', 'high', 'borderline'])).sum())} notable values",
        "CSV + JSON summary included",
    ]
    if test_date:
        stat_lines.insert(0, f"Test date: {test_date}")
    for idx, line in enumerate(stat_lines):
        p = stf.paragraphs[0] if idx == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    total = len(df)
    abnormal = int((df["status"].isin(["low", "high", "borderline"])).sum())
    info = slide.shapes.add_textbox(Inches(4.05), Inches(3.78), Inches(3.8), Inches(0.9))
    itf = info.text_frame
    itf.clear()
    p = itf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"{total} extracted markers"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = itf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = f"{abnormal} notable values"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEAL


def add_overview_boxes(slide, df: pd.DataFrame) -> None:
    add_modern_background(slide)
    add_slide_title(slide, "Executive Overview", "Clean summary of the extracted blood test data")
    metrics = [
        ("Total markers", str(len(df))),
        ("Normal", str(int((df["status"] == "normal").sum()))),
        ("Notable", str(int((df["status"].isin(["low", "high", "borderline"])).sum()))),
    ]
    left = Inches(0.65)
    top = Inches(1.65)
    width = Inches(2.9)
    height = Inches(1.25)

    for idx, (label, value) in enumerate(metrics):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + Inches(idx * 3.05),
            top,
            width,
            height,
        )
        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = LINE
        text_frame = box.text_frame
        text_frame.clear()
        p1 = text_frame.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run1 = p1.add_run()
        run1.text = label
        run1.font.size = Pt(14)
        run1.font.bold = True
        run1.font.color.rgb = SLATE
        p2 = text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = value
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = NAVY


def add_table_slide(prs: Presentation, df: pd.DataFrame, title: str, rows_limit: int = 10) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    add_slide_title(slide, title, "Top extracted values from the OCR output")

    rows = min(len(df), rows_limit) + 1
    table = slide.shapes.add_table(rows, 5, Inches(0.4), Inches(1.55), Inches(9.2), Inches(4.6)).table
    headers = ["Marker", "Value", "Unit", "Reference", "Status"]
    widths = [2.5, 1.1, 1.0, 2.8, 1.2]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
        cell = table.cell(0, idx)
        cell.text = headers[idx]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = NAVY

    for row_idx, (_, row) in enumerate(df.head(rows_limit).iterrows(), start=1):
        values = [
            row["marker"],
            "" if pd.isna(row["value"]) else str(row["value"]),
            row["unit"],
            row["reference"],
            row["status"],
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            if col_idx == 4:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(100, 100, 100)
                if value == "normal":
                    fill.fore_color.rgb = RGBColor(46, 125, 50)
                elif value == "borderline":
                    fill.fore_color.rgb = RGBColor(249, 168, 37)
                elif value == "low":
                    fill.fore_color.rgb = RGBColor(239, 108, 0)
                elif value == "high":
                    fill.fore_color.rgb = RGBColor(198, 40, 40)
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)


def add_alert_slide(prs: Presentation, df: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    add_slide_title(slide, "Abnormal Values", "Values outside or close to the reference range")

    notable = df[df["status"].isin(["low", "high", "borderline"])].copy()
    if notable.empty:
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.5), Inches(1.2))
        tf = text_box.text_frame
        tf.text = "No notable values were detected in the structured data."
        tf.paragraphs[0].font.size = Pt(22)
        return

    notable = notable.sort_values(by=["status", "marker"])
    rows = min(len(notable), 8) + 1
    table = slide.shapes.add_table(rows, 5, Inches(0.4), Inches(1.55), Inches(9.2), Inches(4.6)).table
    headers = ["Marker", "Value", "Unit", "Reference", "Status"]
    widths = [2.5, 1.1, 1.0, 2.8, 1.2]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
        cell = table.cell(0, idx)
        cell.text = headers[idx]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = NAVY

    for row_idx, (_, row) in enumerate(notable.head(8).iterrows(), start=1):
        values = [
            row["marker"],
            "" if pd.isna(row["value"]) else str(row["value"]),
            row["unit"],
            row["reference"],
            row["status"],
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            if col_idx == 4:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(198, 40, 40) if value == "high" else RGBColor(239, 108, 0) if value == "low" else RGBColor(249, 168, 37)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.font.bold = True


def add_chart_slide(prs: Presentation, bar_path: Optional[Path], pie_path: Optional[Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    add_slide_title(slide, "Visual Summary", "Simple charts derived from the extracted values")

    if bar_path and bar_path.exists():
        slide.shapes.add_picture(str(bar_path), Inches(0.3), Inches(1.6), width=Inches(4.85))
    if pie_path and pie_path.exists():
        slide.shapes.add_picture(str(pie_path), Inches(5.1), Inches(1.75), width=Inches(3.65))

    if not bar_path and not pie_path:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tf = text_box.text_frame
        tf.text = "No charts could be created from the available values."
        tf.paragraphs[0].font.size = Pt(20)


def add_summary_slide(prs: Presentation, df: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    add_slide_title(slide, "Normal Values", "Main findings from the structured OCR output")

    body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(8.4), Inches(4.4))
    body = body_box.text_frame
    body.clear()

    total = len(df)
    status_counts = df["status"].value_counts().to_dict()
    notable = df[df["status"].isin(["low", "high", "borderline"])].copy()

    normal_count = int(status_counts.get('normal', 0))
    notable_count = int(status_counts.get('borderline', 0)) + int(status_counts.get('low', 0)) + int(status_counts.get('high', 0))

    bullets = [f"Total markers evaluated: {total}"]
    bullets.append(f"Normal values: {normal_count}")
    bullets.append(f"Notable values: {notable_count}")

    if not notable.empty:
        top_row = notable.copy()
        top_row["abs_dev"] = top_row["pct_deviation_from_center"].abs().fillna(0)
        top_row = top_row.sort_values(by="abs_dev", ascending=False).iloc[0]
        bullets.append(f"Largest deviation: {top_row['marker']} ({top_row['value']} {top_row['unit']})")
    else:
        bullets.append("No notable deviations were detected in the structured values.")

    bullets.append("Use this report together with the original laboratory context for interpretation.")

    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = SLATE


def build_recommendations(df: pd.DataFrame) -> list[str]:
    recommendations = []
    notable = df[df["status"].isin(["low", "high", "borderline"])].copy()

    if notable.empty:
        return [
            "No clear abnormalities detected in the structured data.",
            "Continue with routine follow-up according to the treating clinician's plan.",
            "Use the original laboratory report and symptoms for final interpretation.",
        ]

    lower_markers = " ".join(notable["marker"].str.lower().tolist())

    if any(key in lower_markers for key in ["cholesterin", "hdl", "ldl", "triglyceride"]):
        recommendations.append("Review lipid values with your clinician, especially if they remain above target on repeat testing.")
    if any(key in lower_markers for key in ["glukose", "glucose", "hba1c", "hba1c"]):
        recommendations.append("Correlate glucose-related values with fasting status and diabetes risk factors.")
    if any(key in lower_markers for key in ["kreatinin", "egfr", "harnstoff", "harnsäure"]):
        recommendations.append("Discuss kidney-related markers in clinical context and consider follow-up if values stay abnormal.")
    if any(key in lower_markers for key in ["gpt", "alat", "got", "ast", "gamma-gt", "bilirubin", "alkalische phosphatase"]):
        recommendations.append("Correlate liver-related markers with medications, alcohol intake, and recent illness before interpreting further.")
    if any(key in lower_markers for key in ["crp", "hs-crp", "leukozyten"]):
        recommendations.append("Inflammation-related markers may warrant clinical correlation if symptoms are present.")
    if any(key in lower_markers for key in ["ferritin", "eisen", "transferrin", "vitamin b12", "folsäure"]):
        recommendations.append("Consider iron and vitamin status follow-up if fatigue, anemia, or dietary concerns are present.")
    if any(key in lower_markers for key in ["tsh", "ft3", "ft4"]):
        recommendations.append("If thyroid markers are abnormal, repeat testing and symptom review may be appropriate.")

    if not recommendations:
        recommendations.append("Review the notable values with the treating clinician for clinical context.")

    recommendations.append("This report is informational and does not replace medical advice.")
    return recommendations


def add_recommendations_slide(prs: Presentation, df: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    add_slide_title(slide, "Recommendations", "Data-driven next steps based on notable values")

    recommendations = build_recommendations(df)
    body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(8.6), Inches(4.6))
    body = body_box.text_frame
    body.clear()

    for idx, item in enumerate(recommendations):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = SLATE


def build_word_report(df: pd.DataFrame, test_date: Optional[str]) -> None:
    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = DocxPt(10.5)

    title = doc.add_paragraph()
    title_run = title.add_run("Doctor's Report")
    title_run.bold = True
    title_run.font.size = DocxPt(18)

    if test_date:
        p = doc.add_paragraph()
        run = p.add_run(f"Blood test date: {test_date}")
        run.bold = True

    doc.add_paragraph(f"Total markers extracted: {len(df)}")
    doc.add_paragraph(f"Normal values: {int((df['status'] == 'normal').sum())}")
    doc.add_paragraph(f"Abnormal values: {int((df['status'].isin(['low', 'high', 'borderline'])).sum())}")

    normal_df = df[df["status"] == "normal"].sort_values(by="marker")
    abnormal_df = df[df["status"].isin(["low", "high", "borderline"])].sort_values(by=["status", "marker"])

    def add_section_heading(text: str) -> None:
        heading = doc.add_paragraph()
        run = heading.add_run(text)
        run.bold = True
        run.font.size = DocxPt(13)

    def add_table(section_df: pd.DataFrame) -> None:
        if section_df.empty:
            doc.add_paragraph("No values found in this section.")
            return
        table = doc.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        headers = ["Marker", "Value", "Unit", "Reference"]
        for idx, header in enumerate(headers):
            table.rows[0].cells[idx].text = header
        for _, row in section_df.iterrows():
            cells = table.add_row().cells
            cells[0].text = str(row["marker"])
            cells[1].text = "" if pd.isna(row["value"]) else str(row["value"])
            cells[2].text = str(row["unit"])
            cells[3].text = str(row["reference"])

    add_section_heading("Normal Values")
    add_table(normal_df)
    add_section_heading("Abnormal Values")
    add_table(abnormal_df)

    add_section_heading("Recommendations")
    for item in build_recommendations(df):
        doc.add_paragraph(item, style="List Bullet")

    doc.save(DOCX_PATH)


def add_compact_table_slide(prs: Presentation, df: pd.DataFrame, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide)
    add_slide_title(slide, title, subtitle)

    if df.empty:
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.5), Inches(1.2))
        tf = text_box.text_frame
        tf.text = "No values matched this section."
        tf.paragraphs[0].font.size = Pt(22)
        return

    rows = min(len(df), 10) + 1
    table = slide.shapes.add_table(rows, 4, Inches(0.4), Inches(1.55), Inches(9.2), Inches(4.6)).table
    headers = ["Marker", "Value", "Unit", "Reference"]
    widths = [3.2, 1.3, 1.2, 3.5]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
        cell = table.cell(0, idx)
        cell.text = headers[idx]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = NAVY

    for row_idx, (_, row) in enumerate(df.head(10).iterrows(), start=1):
        values = [
            row["marker"],
            "" if pd.isna(row["value"]) else str(row["value"]),
            row["unit"],
            row["reference"],
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)


def main() -> None:
    try:
        ensure_output_dir(OUTPUT_DIR)
        df_raw = load_records(CSV_PATH, JSON_PATH)
        df = enrich_dataframe(df_raw)
        metadata = load_metadata(METADATA_PATH)
        test_date = metadata.get("test_date")

        bar_chart_path = build_bar_chart(df, OUTPUT_DIR / "ppt_bar_chart.png")
        deviation_chart_path = build_deviation_chart(df, OUTPUT_DIR / "ppt_deviation_chart.png")
        pie_chart_path = build_status_pie(df, OUTPUT_DIR / "ppt_status_pie.png")

        prs = Presentation()
        add_title_slide(prs, df, test_date)

        overview = prs.slides.add_slide(prs.slide_layouts[6])
        add_overview_boxes(overview, df)

        normal_df = df[df["status"] == "normal"].sort_values(by="marker")
        abnormal_df = df[df["status"].isin(["low", "high", "borderline"])].sort_values(by=["status", "marker"])

        add_compact_table_slide(prs, normal_df, "Normal Values", "Values inside the reference range")
        add_compact_table_slide(prs, abnormal_df, "Abnormal Values", "Values outside or close to the reference range")
        add_alert_slide(prs, df)
        add_chart_slide(prs, bar_chart_path, pie_chart_path)
        add_recommendations_slide(prs, df)

        prs.save(PPTX_PATH)
        build_word_report(df, test_date)

        print(f"Saved PowerPoint: {PPTX_PATH}")
        print(f"Saved Word report: {DOCX_PATH}")
        print(f"Loaded data rows: {len(df)}")
        print(f"Notable rows: {int((df['status'].isin(['low', 'high', 'borderline'])).sum())}")
    except Exception as exc:
        print(f"Failed to generate presentation: {exc}")
        sys.exit(1)


if __name__ == "__main__":
    main()
